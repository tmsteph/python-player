#this is the main file we will be working with. 

#the code below prints hello world to console

#type the command "python test.py" without the question marks in the console to run the program

from pptx import Presentation
from pptx.util import Inches
prs = Presentation()

title_slide_layout = prs.slide_layouts[0]

slide = prs.slides.add_slide(title_slide_layout)

title = slide.shapes.title

subtitle = slide.placeholders[1]

title.text = "Test Slide"
subtitle.text = "tmsteph was here!"

prs.save('test.pptx')

print("Saved Powerpoint")

#def pptx_to_pdf(pptx_path, pdf_path):
    # Load the PowerPoint file
#    prs1 = Presentation(pptx_path)

    #Save the PowerPoint file as  PDF file
#    prs1.save(pdf_path)

# example usage
#pptx_to_pdf('test.pptx','test.pdf')

print("changed to pdf")

#libreoffice --headless --convert-to pdf ~/python-practice/python-player/test.pptx

import os
cmd = 'ls'
os.system(cmd)


import subprocess

input_file = '/path/to/test.pptx'
output_dir = '/path/to/output/dir'

subprocess.call(['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', output_dir, input_file])
